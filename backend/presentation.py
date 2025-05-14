from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import nltk
from nltk.tokenize import sent_tokenize
import os
import fitz  # PyMuPDF
from PIL import Image, ImageDraw, ImageFont

# Download required NLTK data (run this once)
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')
    nltk.download('punkt_tab')

def create_presentation(text, output_path, is_formal=True):
    """
    Create a presentation from the summarized text
    Args:
        text: The summarized text to use in the presentation
        output_path: Where to save the presentation
        is_formal: Whether to create a formal or informal presentation
    """
    if is_formal:
        return generate_formal_presentation(text, output_path)
    else:
        return generate_informal_presentation(text, output_path)

def classify_text(text, keywords):
    """Classify sentences into sections based on keywords"""
    classified_data = {section: [] for section in keywords}
    sentences = sent_tokenize(text)
    
    for sentence in sentences:
        sentence = sentence.strip()
        for section, words in keywords.items():
            if any(word.lower() in sentence.lower() for word in words):
                classified_data[section].append(sentence)
                break
            
    # If a sentence doesn't match any section, put it in the most relevant section
    for sentence in sentences:
        sentence = sentence.strip()
        found = False
        for section in classified_data:
            if sentence in classified_data[section]:
                found = True
                break
        
        if not found:
            if any(word in sentence.lower() for word in ['found', 'showed', 'revealed', 'indicates']):
                classified_data['Results'].append(sentence)
            else:
                classified_data['Methodology'].append(sentence)
    
    return classified_data

def generate_formal_presentation(text, output_path):
    """Generate a formal presentation"""
    # Define keywords for each section
    section_keywords = {
        "Abstract": ["abstract", "summary", "overview"],
        "Introduction": ["introduction", "background", "context"],
        "Methodology": ["methodology", "method", "approach", "process"],
        "Results": ["result", "finding", "observation", "show", "reveal"],
        "Conclusion": ["conclusion", "summary", "finally", "overall"]
    }

    # Classify the text into sections
    classified_sections = classify_text(text, section_keywords)

    # Create presentation
    prs = Presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Research Summary"
    slide.placeholders[1].text = "Formal Presentation"

    # Add content slides
    for section, content in classified_sections.items():
        if content:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = section
            
            text_frame = slide.shapes.placeholders[1].text_frame
            text_frame.clear()

            for sentence in content:
                p = text_frame.add_paragraph()
                p.text = "• " + sentence
                p.font.size = Pt(14)
                p.space_after = Pt(12)

    # Add conclusion slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Thank You"
    text_frame = slide.shapes.placeholders[1].text_frame
    p = text_frame.add_paragraph()
    p.text = "Thank you for your attention!"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save(output_path)
    return output_path

def generate_informal_presentation(text, output_path):
    """Generate an informal presentation"""
    # Create presentation
    prs = Presentation()

    # Add title slide with fun title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Quick Summary!"
    subtitle.text = "The Fun Version 😊"
    
    # Split text into sentences for key points
    sentences = sent_tokenize(text)
    
    # Create content slides with bullet points
    for i in range(0, len(sentences), 3):  # 3 points per slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        
        title.text = f"Key Points {i//3 + 1}"
        tf = body.text_frame
        
        for j in range(3):
            if i + j < len(sentences):
                p = tf.add_paragraph()
                p.text = "• " + sentences[i + j]
                p.font.size = Pt(18)
                p.space_after = Pt(12)

    # Add fun closing slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thanks for watching! 🎉"
    p.font.size = Pt(40)
    p.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save(output_path)
    return output_path
